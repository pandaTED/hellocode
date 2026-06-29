"""Knowledge base engine — file indexing, text extraction, and search."""

from __future__ import annotations

import csv
import hashlib
import json
import logging
import re
from io import StringIO
from pathlib import Path

from .storage import Storage

logger = logging.getLogger("hellocode.knowledge")

SUPPORTED_TYPES = {
    "md", "txt", "pdf", "docx", "doc", "xlsx", "xls",
    "pptx", "ppt", "csv", "json", "yaml", "yml",
}


def _file_hash(fp: Path) -> str:
    try:
        stat = fp.stat()
    except OSError:
        return ""
    h = hashlib.md5()
    h.update(str(stat.st_size).encode())
    h.update(str(int(stat.st_mtime * 1000)).encode())
    return h.hexdigest()


def _chunk_text(text: str, chunk_size: int = 1000, overlap: int = 100) -> list[dict]:
    if not text or not text.strip():
        return []
    if overlap >= chunk_size:
        overlap = chunk_size // 4
    chunks = []
    start = 0
    idx = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append({
                "index": idx,
                "content": chunk.strip(),
                "offset": start,
            })
            idx += 1
        start = end - overlap if end < len(text) else end
    return chunks


class KnowledgeEngine:
    def __init__(self, storage: Storage, data_dir: Path, chunk_size: int = 1000,
                 chunk_overlap: int = 100, max_file_size_mb: int = 50):
        self.storage = storage
        self.data_dir = data_dir / "knowledge"
        self.data_dir.mkdir(parents=True, exist_ok=True)
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.max_file_size_mb = max_file_size_mb

    # ── Source Management ──

    def add_source(self, name: str, path: Path) -> dict:
        path = path.resolve()
        if not path.exists():
            raise FileNotFoundError(f"Path not found: {path}")
        existing = self.storage.get_kb_source_by_path(str(path))
        if existing:
            return existing
        source_type = "folder" if path.is_dir() else "file"
        source = self.storage.create_kb_source(
            id=self.storage.uid(), name=name, path=str(path), source_type=source_type,
        )
        return source

    def remove_source(self, source_id: str) -> None:
        self.storage.delete_kb_source(source_id)

    def list_sources(self) -> list[dict]:
        return self.storage.list_kb_sources()

    def get_source(self, source_id: str) -> dict | None:
        return self.storage.get_kb_source(source_id)

    # ── Indexing ──

    def index_source(self, source_id: str) -> dict:
        source = self.storage.get_kb_source(source_id)
        if not source:
            raise ValueError(f"Source not found: {source_id}")
        path = Path(source["path"])
        if not path.exists():
            self.storage.update_kb_source(source_id, status="error")
            return {"error": f"Path not found: {path}"}

        files = []
        max_size = self.max_file_size_mb * 1024 * 1024
        if path.is_dir():
            for fp in sorted(path.rglob("*")):
                if fp.is_file() and fp.suffix.lstrip(".").lower() in SUPPORTED_TYPES:
                    try:
                        if fp.stat().st_size <= max_size:
                            files.append(fp)
                    except OSError:
                        continue
        else:
            if path.suffix.lstrip(".").lower() in SUPPORTED_TYPES:
                files.append(path)

        indexed, skipped, errors = 0, 0, 0
        for fp in files:
            result = self._index_file(fp, source_id)
            if result == "indexed":
                indexed += 1
            elif result == "skipped":
                skipped += 1
            else:
                errors += 1

        self.storage.update_kb_source(
            source_id, file_count=indexed + skipped + errors,
            last_indexed_at=self.storage.now(),
        )
        self.storage.rebuild_kb_chunk_fts()
        return {"indexed": indexed, "skipped": skipped, "errors": errors, "total": len(files)}

    def _index_file(self, file_path: Path, source_id: str) -> str:
        content_hash = _file_hash(file_path)
        existing = self.storage.get_kb_document_by_path(str(file_path))
        if existing and existing.get("content_hash") == content_hash and existing.get("status") == "indexed":
            return "skipped"

        file_type = file_path.suffix.lstrip(".").lower()
        file_name = file_path.name
        file_size = file_path.stat().st_size

        if existing:
            doc_id = existing["id"]
            self.storage.update_kb_document(doc_id, status="pending", content_hash=content_hash)
        else:
            doc_id = self.storage.uid()
            self.storage.create_kb_document(
                id=doc_id, source_id=source_id, file_path=str(file_path),
                file_name=file_name, file_type=file_type, file_size=file_size,
                content_hash=content_hash,
            )

        try:
            pages = self._extract_text(file_path, file_type)
            full_text = "\n\n".join(p["text"] for p in pages if p["text"].strip())
            metadata = json.dumps({"pages": len(pages)}, ensure_ascii=False)
            self.storage.update_kb_document(
                doc_id, extracted_text=full_text, metadata=metadata, status="indexed",
            )
            self._create_chunks(doc_id, pages)
            return "indexed"
        except Exception as e:
            logger.error("Failed to index %s: %s", file_path, e)
            self.storage.update_kb_document(doc_id, status="error", error_message=str(e))
            return "error"

    def _create_chunks(self, doc_id: str, pages: list[dict]) -> None:
        self.storage._delete_kb_chunks_for_doc(doc_id)
        chunk_data = []
        fts_entries = []
        chunk_idx = 0
        for page in pages:
            text = page["text"]
            page_num = page.get("page_number")
            heading = page.get("heading")
            chunks = _chunk_text(text, self.chunk_size, self.chunk_overlap)
            for c in chunks:
                chunk_id = self.storage.uid()
                chunk_data.append({
                    "id": chunk_id, "document_id": doc_id,
                    "chunk_index": chunk_idx, "content": c["content"],
                    "page_number": page_num, "heading": heading, "char_offset": c["offset"],
                })
                fts_entries.append((chunk_id, c["content"], heading or ""))
                chunk_idx += 1
        if chunk_data:
            self.storage.create_kb_chunks_batch(chunk_data)
            self.storage.index_kb_chunks_fts_batch(fts_entries)

    def _extract_text(self, file_path: Path, file_type: str) -> list[dict]:
        if file_type in ("md", "txt"):
            return self._extract_plain(file_path)
        elif file_type == "pdf":
            return self._extract_pdf(file_path)
        elif file_type in ("docx", "doc"):
            return self._extract_docx(file_path)
        elif file_type in ("xlsx", "xls"):
            return self._extract_xlsx(file_path)
        elif file_type in ("pptx", "ppt"):
            return self._extract_pptx(file_path)
        elif file_type == "csv":
            return self._extract_csv(file_path)
        elif file_type in ("json", "yaml", "yml"):
            return self._extract_data_file(file_path)
        return [{"text": "", "page_number": None, "heading": None}]

    def _extract_plain(self, fp: Path) -> list[dict]:
        text = fp.read_text(encoding="utf-8", errors="replace")
        return [{"text": text, "page_number": None, "heading": None}]

    def _extract_pdf(self, fp: Path) -> list[dict]:
        try:
            import fitz
            doc = fitz.open(str(fp))
            pages = []
            for i, page in enumerate(doc):
                text = page.get_text()
                heading = self._detect_heading(text)
                pages.append({"text": text, "page_number": i + 1, "heading": heading})
            doc.close()
            return pages
        except ImportError:
            logger.warning("PyMuPDF not installed, skipping PDF: %s", fp)
            return [{"text": fp.stem, "page_number": None, "heading": None}]
        except Exception as e:
            logger.error("PDF parse error %s: %s", fp, e)
            return [{"text": "", "page_number": None, "heading": None}]

    def _extract_docx(self, fp: Path) -> list[dict]:
        try:
            from docx import Document
            doc = Document(str(fp))
            paragraphs = []
            current_heading = None
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                if para.style and para.style.name.startswith("Heading"):
                    current_heading = text
                paragraphs.append(text)
            full_text = "\n\n".join(paragraphs)
            heading = self._detect_heading(full_text)
            return [{"text": full_text, "page_number": None, "heading": heading}]
        except ImportError:
            logger.warning("python-docx not installed, skipping: %s", fp)
            return [{"text": fp.stem, "page_number": None, "heading": None}]

    def _extract_xlsx(self, fp: Path) -> list[dict]:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(fp), read_only=True, data_only=True)
            pages = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    text = f"Sheet: {sheet_name}\n" + "\n".join(rows)
                    pages.append({"text": text, "page_number": None, "heading": f"Sheet: {sheet_name}"})
            wb.close()
            return pages or [{"text": fp.stem, "page_number": None, "heading": None}]
        except ImportError:
            logger.warning("openpyxl not installed, skipping: %s", fp)
            return [{"text": fp.stem, "page_number": None, "heading": None}]

    def _extract_pptx(self, fp: Path) -> list[dict]:
        try:
            from pptx import Presentation
            prs = Presentation(str(fp))
            pages = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    text = "\n".join(texts)
                    heading = self._detect_heading(text)
                    pages.append({"text": text, "page_number": i + 1, "heading": heading})
            return pages or [{"text": fp.stem, "page_number": None, "heading": None}]
        except ImportError:
            logger.warning("python-pptx not installed, skipping: %s", fp)
            return [{"text": fp.stem, "page_number": None, "heading": None}]

    def _extract_csv(self, fp: Path) -> list[dict]:
        try:
            text = fp.read_text(encoding="utf-8", errors="replace")
            reader = csv.reader(StringIO(text))
            rows = list(reader)
            if not rows:
                return [{"text": "", "page_number": None, "heading": None}]
            header = " | ".join(rows[0])
            body = "\n".join(" | ".join(row) for row in rows[1:])
            full_text = f"CSV Columns: {header}\n{body}"
            return [{"text": full_text, "page_number": None, "heading": None}]
        except Exception as e:
            logger.error("CSV parse error %s: %s", fp, e)
            return [{"text": "", "page_number": None, "heading": None}]

    def _extract_data_file(self, fp: Path) -> list[dict]:
        text = ""
        try:
            text = fp.read_text(encoding="utf-8", errors="replace")
            if fp.suffix in (".json",):
                data = json.loads(text)
                formatted = json.dumps(data, ensure_ascii=False, indent=2)
            else:
                formatted = text
            return [{"text": formatted, "page_number": None, "heading": None}]
        except Exception as e:
            logger.error("Data file parse error %s: %s", fp, e)
            return [{"text": text, "page_number": None, "heading": None}]

    @staticmethod
    def _detect_heading(text: str) -> str | None:
        for line in text.split("\n")[:10]:
            line = line.strip()
            if not line:
                continue
            m = re.match(r"^(#{1,6})\s+(.+)", line)
            if m:
                return m.group(2).strip()
            if len(line) < 100 and (line[0].isupper() or line[0].isdigit()):
                return line
            break
        return None

    # ── Search ──

    def search(self, query: str, source_id: str | None = None,
               file_type: str | None = None, limit: int = 10) -> list[dict]:
        return self.storage.search_kb_chunks(query, source_id, file_type, limit)

    # ── Document Management ──

    def get_document(self, doc_id: str) -> dict | None:
        return self.storage.get_kb_document(doc_id)

    def get_document_chunks(self, doc_id: str) -> list[dict]:
        return self.storage._execute(
            "SELECT * FROM kb_chunk WHERE document_id=? ORDER BY chunk_index", (doc_id,)
        )

    def list_documents(self, source_id: str | None = None) -> list[dict]:
        return self.storage.list_kb_documents(source_id)

    def get_stats(self) -> dict:
        return self.storage.get_kb_stats()
